#!/usr/bin/env python3
"""
PPTX to Markdown Converter

A tool that converts PowerPoint (.pptx) files to Markdown format,
preserving text, structure, and extracting images.
"""

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("Error: python-pptx is required. Install it with: pip install python-pptx")
    sys.exit(1)


def extract_text_from_shape(shape) -> str:
    """Extract text from a shape, handling different shape types."""
    text_parts = []

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            para_text = ""
            for run in paragraph.runs:
                text = run.text
                if text.strip():
                    # Check for bold/italic
                    is_bold = run.font.bold
                    is_italic = run.font.italic

                    if is_bold and is_italic:
                        text = f"***{text}***"
                    elif is_bold:
                        text = f"**{text}**"
                    elif is_italic:
                        text = f"*{text}*"

                    para_text += text
                else:
                    para_text += text

            if para_text.strip():
                # Detect bullet level
                level = paragraph.level if paragraph.level else 0
                indent = "  " * level

                # Check if it's a bullet point
                if level > 0 or (paragraph.text.strip() and not para_text.startswith("#")):
                    text_parts.append(f"{indent}- {para_text.strip()}")
                else:
                    text_parts.append(para_text.strip())

    return "\n".join(text_parts)


def extract_table(shape) -> str:
    """Extract table content as Markdown table."""
    if not shape.has_table:
        return ""

    table = shape.table
    rows = []

    for row_idx, row in enumerate(table.rows):
        cells = []
        for cell in row.cells:
            cell_text = cell.text.strip().replace("|", "\\|")
            cells.append(cell_text)
        rows.append("| " + " | ".join(cells) + " |")

        # Add header separator after first row
        if row_idx == 0:
            separator = "| " + " | ".join(["---"] * len(cells)) + " |"
            rows.append(separator)

    return "\n".join(rows)


def extract_image(shape, slide_num: int, img_index: int, output_dir: Path) -> str:
    """Extract image from shape and save it."""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            image_bytes = image.blob
            image_ext = image.ext

            # Create images directory
            images_dir = output_dir / "images"
            images_dir.mkdir(exist_ok=True)

            # Save image
            image_filename = f"slide{slide_num}_img{img_index}.{image_ext}"
            image_path = images_dir / image_filename

            with open(image_path, "wb") as f:
                f.write(image_bytes)

            return f"![Image {img_index}](images/{image_filename})"
    except Exception as e:
        print(f"Warning: Could not extract image from slide {slide_num}: {e}")

    return ""


def get_shape_sort_key(shape):
    """Get sort key for shapes based on position (top to bottom, left to right)."""
    try:
        top = shape.top if shape.top else 0
        left = shape.left if shape.left else 0
        return (top, left)
    except Exception:
        return (0, 0)


def convert_slide_to_markdown(slide, slide_num: int, output_dir: Path,
                               extract_images: bool) -> str:
    """Convert a single slide to Markdown."""
    markdown_parts = []
    img_index = 0

    # Sort shapes by position (top to bottom, left to right)
    shapes = sorted(slide.shapes, key=get_shape_sort_key)

    title_found = False

    for shape in shapes:
        # Handle title
        if shape.is_placeholder and hasattr(shape, 'placeholder_format'):
            ph_type = shape.placeholder_format.type
            # Title placeholder types: TITLE (1), CENTER_TITLE (3), SUBTITLE (4)
            if ph_type in [1, 3] and not title_found:
                title_text = shape.text.strip()
                if title_text:
                    markdown_parts.append(f"## {title_text}")
                    title_found = True
                continue
            elif ph_type == 4:  # Subtitle
                subtitle_text = shape.text.strip()
                if subtitle_text:
                    markdown_parts.append(f"*{subtitle_text}*")
                continue

        # Handle images
        if extract_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img_index += 1
            img_md = extract_image(shape, slide_num, img_index, output_dir)
            if img_md:
                markdown_parts.append(img_md)
            continue

        # Handle tables
        if shape.has_table:
            table_md = extract_table(shape)
            if table_md:
                markdown_parts.append(table_md)
            continue

        # Handle text
        if shape.has_text_frame:
            text = extract_text_from_shape(shape)
            if text.strip():
                markdown_parts.append(text)

    return "\n\n".join(markdown_parts)


def pptx_to_markdown(pptx_path: str, output_path: str = None,
                     extract_images: bool = True,
                     slide_range: tuple = None) -> str:
    """
    Convert a PowerPoint file to Markdown format.

    Args:
        pptx_path: Path to the input PPTX file
        output_path: Path for the output Markdown file (optional)
        extract_images: Whether to extract and include images
        slide_range: Tuple of (start_slide, end_slide) for partial conversion (1-indexed)

    Returns:
        The generated Markdown content
    """
    pptx_path = Path(pptx_path)

    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    # Determine output path
    if output_path:
        output_path = Path(output_path)
    else:
        output_path = pptx_path.with_suffix('.md')

    output_dir = output_path.parent
    output_dir.mkdir(parents=True, exist_ok=True)

    print(f"Converting: {pptx_path}")
    print(f"Output: {output_path}")

    # Open presentation
    prs = Presentation(pptx_path)
    total_slides = len(prs.slides)

    print(f"Total slides: {total_slides}")

    # Determine slide range
    start_slide = 0
    end_slide = total_slides

    if slide_range:
        start_slide = max(0, slide_range[0] - 1)
        end_slide = min(total_slides, slide_range[1])

    # Convert slides
    markdown_content = []
    markdown_content.append(f"# {pptx_path.stem}\n")
    markdown_content.append(f"*Converted from PowerPoint - {total_slides} slides*\n")
    markdown_content.append("---\n")

    for slide_num in range(start_slide, end_slide):
        print(f"Processing slide {slide_num + 1}/{end_slide}...", end='\r')

        slide = prs.slides[slide_num]
        slide_md = convert_slide_to_markdown(
            slide, slide_num + 1, output_dir, extract_images
        )

        markdown_content.append(f"\n---\n\n### Slide {slide_num + 1}\n")

        if slide_md.strip():
            markdown_content.append(slide_md)
        else:
            markdown_content.append("*(Empty slide)*")

    print(f"\nProcessed {end_slide - start_slide} slides.")

    # Join content
    final_content = "\n\n".join(markdown_content)

    # Write output
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(final_content)

    print(f"Markdown saved to: {output_path}")

    return final_content


def main():
    parser = argparse.ArgumentParser(
        description="Convert PowerPoint files to Markdown format",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python pptx_to_markdown.py presentation.pptx
  python pptx_to_markdown.py presentation.pptx -o output.md
  python pptx_to_markdown.py presentation.pptx --slides 1-10
  python pptx_to_markdown.py presentation.pptx --no-images
        """
    )

    parser.add_argument("pptx_file", help="Path to the PPTX file to convert")
    parser.add_argument("-o", "--output", help="Output Markdown file path")
    parser.add_argument("--no-images", action="store_true",
                        help="Skip image extraction")
    parser.add_argument("--slides", help="Slide range to convert (e.g., '1-10')")

    args = parser.parse_args()

    # Parse slide range
    slide_range = None
    if args.slides:
        try:
            parts = args.slides.split('-')
            if len(parts) == 2:
                slide_range = (int(parts[0]), int(parts[1]))
            else:
                slide_range = (int(parts[0]), int(parts[0]))
        except ValueError:
            print(f"Invalid slide range: {args.slides}")
            sys.exit(1)

    try:
        pptx_to_markdown(
            args.pptx_file,
            output_path=args.output,
            extract_images=not args.no_images,
            slide_range=slide_range
        )
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
