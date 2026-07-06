"""Regression test for pptx -> markdown bullet handling.

Bug: every non-empty, non-heading paragraph was rendered as a Markdown bullet,
so ordinary body text became a bogus list. A plain (unbulleted) paragraph must
now render as plain text.
"""

import importlib.util
from pathlib import Path

import pytest

pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

MODULE = Path(__file__).resolve().parents[1] / "scripts" / "pptx_to_markdown.py"
spec = importlib.util.spec_from_file_location("pptx_to_markdown", MODULE)
mod = importlib.util.module_from_spec(spec)
spec.loader.exec_module(mod)


def test_plain_textbox_paragraph_is_not_bulleted():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
    tf = box.text_frame
    tf.text = "Just a plain sentence."
    p2 = tf.add_paragraph()
    p2.text = "Another plain line."

    md = mod.extract_text_from_shape(box)
    lines = [l for l in md.splitlines() if l.strip()]
    assert lines == ["Just a plain sentence.", "Another plain line."]
    assert not any(l.lstrip().startswith("- ") for l in lines)


def test_is_bulleted_false_for_plain_paragraph():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    box.text_frame.text = "x"
    para = box.text_frame.paragraphs[0]
    assert mod._is_bulleted(para) is False
